"""
Generates the Spring Music Modernization Report PPTX.
Run with: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ────────────────────────────────────────────────────────────
SPRING_GREEN   = RGBColor(0x6D, 0xB3, 0x3F)   # Spring green
DARK_GRAY      = RGBColor(0x2D, 0x2D, 0x2D)   # slide background / headings
MID_GRAY       = RGBColor(0x55, 0x55, 0x55)   # body text
LIGHT_GRAY     = RGBColor(0xF4, 0xF4, 0xF4)   # card/table fill
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE    = RGBColor(0x1E, 0x6B, 0xB0)
ACCENT_ORANGE  = RGBColor(0xE0, 0x7B, 0x20)
ACCENT_RED     = RGBColor(0xC0, 0x39, 0x2B)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

BLANK = prs.slide_layouts[6]   # completely blank


# ═══════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(14), bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_bullet_box(slide, items, l, t, w, h,
                   size=Pt(13), color=MID_GRAY, bullet="•  ", line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = size
        run.font.color.rgb = color
    return tb


def slide_bg(slide, color=DARK_GRAY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)


def green_bar(slide, title_text, subtitle_text=None):
    """Top accent bar + title."""
    add_rect(slide, 0, 0, 13.33, 0.08, fill=SPRING_GREEN)
    add_text(slide, title_text,
             0.4, 0.18, 12.5, 0.65,
             size=Pt(28), bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 0.4, 0.82, 12.5, 0.45,
                 size=Pt(15), color=RGBColor(0xBB, 0xBB, 0xBB))


def bottom_bar(slide, text="Spring Music | Modernization Report  ·  2026"):
    add_rect(slide, 0, 7.2, 13.33, 0.3, fill=RGBColor(0x22, 0x22, 0x22))
    add_text(slide, text, 0.3, 7.22, 12.7, 0.26,
             size=Pt(9), color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.LEFT)


def card(slide, l, t, w, h, fill=LIGHT_GRAY, radius=False):
    return add_rect(slide, l, t, w, h, fill=fill)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

# large green left bar
add_rect(sl, 0, 0, 0.35, 7.5, fill=SPRING_GREEN)

# decorative light rect
add_rect(sl, 0.35, 2.8, 12.98, 0.06, fill=SPRING_GREEN)

add_text(sl, "Spring Music", 0.7, 1.1, 11.5, 1.1,
         size=Pt(52), bold=True, color=WHITE)
add_text(sl, "Modernization Report", 0.7, 2.2, 11.5, 0.75,
         size=Pt(30), color=SPRING_GREEN)
add_text(sl,
         "Current State  ·  Why Modernize  ·  Target State  ·  Key Differences",
         0.7, 3.1, 11.5, 0.55,
         size=Pt(16), color=RGBColor(0xAA, 0xAA, 0xAA))
add_text(sl, "March 2026", 0.7, 6.7, 4.0, 0.45,
         size=Pt(12), color=RGBColor(0x88, 0x88, 0x88))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Agenda")
bottom_bar(sl)

sections = [
    ("01", "Current State",      "Architecture, tech stack, and existing services"),
    ("02", "Why Modernize",      "Pain points, risks, and business drivers"),
    ("03", "Target State",       "Decomposed services, upgraded stack"),
    ("04", "Key Differences",    "Side-by-side comparison of As-Is vs To-Be"),
    ("05", "Migration Summary",  "What changed, what stayed the same"),
]

col_x = [0.5, 3.2, 5.9, 8.6, 11.3]
for i, (num, title, desc) in enumerate(sections):
    x = col_x[i]
    card(sl, x, 1.35, 2.4, 5.5, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_text(sl, num, x + 0.15, 1.5, 2.1, 0.7,
             size=Pt(36), bold=True, color=SPRING_GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, title, x + 0.1, 2.3, 2.2, 0.65,
             size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x + 0.1, 3.1, 2.2, 2.5,
             size=Pt(11), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Current State: Architecture overview
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Current State", "Single Spring Boot monolith — everything in one JAR")
bottom_bar(sl)

# Monolith box
add_rect(sl, 0.4, 1.35, 12.5, 5.55, fill=RGBColor(0x3A, 0x3A, 0x3A),
         line=SPRING_GREEN, line_w=Pt(1.5))
add_text(sl, "spring-music-1.0.jar  (Monolith)",
         0.7, 1.45, 8.0, 0.5, size=Pt(13), bold=True, color=SPRING_GREEN)

# Browser box
add_rect(sl, 0.7, 2.05, 2.3, 0.85, fill=RGBColor(0x1E, 0x6B, 0xB0))
add_text(sl, "Browser\n(AngularJS SPA)", 0.75, 2.1, 2.2, 0.75,
         size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow right
add_text(sl, "REST calls", 3.05, 2.2, 1.3, 0.4,
         size=Pt(9), color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)
add_text(sl, "──────►", 3.0, 2.35, 1.4, 0.4,
         size=Pt(12), color=RGBColor(0x99, 0x99, 0x99))

# Controllers
ctrl_y = 2.0
for label, col in [("AlbumController\n/albums/**", 4.5),
                   ("InfoController\n/appinfo", 6.6),
                   ("ErrorController\n/errors/**", 8.7)]:
    add_rect(sl, col, ctrl_y, 1.75, 0.95, fill=RGBColor(0x4A, 0x4A, 0x4A),
             line=ACCENT_BLUE, line_w=Pt(1))
    add_text(sl, label, col + 0.08, ctrl_y + 0.1, 1.6, 0.75,
             size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)

# Arrow down from AlbumController
add_text(sl, "│\n▼", 5.27, 3.0, 0.5, 0.55,
         size=Pt(12), color=RGBColor(0x99, 0x99, 0x99))

# Repository layer label
add_rect(sl, 4.5, 3.55, 4.0, 0.45, fill=RGBColor(0x55, 0x55, 0x55))
add_text(sl, "CrudRepository <interface>  (single active implementation)",
         4.55, 3.58, 3.9, 0.38, size=Pt(10), color=SPRING_GREEN, align=PP_ALIGN.CENTER)

# Three repo boxes
for label, col, clr in [
        ("JpaAlbumRepository\n(default/mysql/postgres)", 0.7, RGBColor(0x4A, 0x80, 0x3A)),
        ("MongoAlbumRepository\n(mongodb profile)", 4.1, RGBColor(0x2E, 0x6B, 0x44)),
        ("RedisAlbumRepository\n(redis profile)", 7.5, RGBColor(0x7B, 0x2A, 0x2A))]:
    add_rect(sl, col, 4.15, 2.9, 0.95, fill=clr)
    add_text(sl, label, col + 0.1, 4.2, 2.7, 0.85,
             size=Pt(10), color=WHITE, align=PP_ALIGN.CENTER)

# DB layer
add_rect(sl, 0.7, 5.3, 12.0, 0.75, fill=RGBColor(0x3A, 0x3A, 0x3A),
         line=RGBColor(0x66, 0x66, 0x66), line_w=Pt(1))
add_text(sl, "Active Database  —  H2  |  MySQL  |  PostgreSQL  |  MongoDB  |  Redis  |  Oracle  |  SQL Server",
         0.9, 5.42, 11.5, 0.5, size=Pt(11), color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

# Profile resolver note
add_rect(sl, 10.5, 2.05, 2.2, 1.55, fill=RGBColor(0x3A, 0x3A, 0x3A),
         line=ACCENT_ORANGE, line_w=Pt(1))
add_text(sl, "SpringApplication\nContextInitializer\n\nDetects CF services\n→ activates profile",
         10.55, 2.1, 2.1, 1.45, size=Pt(9), color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Current State: Tech stack
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Current State", "Technology stack — as of original codebase")
bottom_bar(sl)

stack_items = [
    ("Runtime",       "Java 8",                  ACCENT_RED),
    ("Framework",     "Spring Boot 2.4.0",        ACCENT_RED),
    ("Build tool",    "Gradle 6.7",               ACCENT_RED),
    ("Testing",       "JUnit 4",                  ACCENT_RED),
    ("Namespace",     "javax.* (Jakarta EE 8)",   ACCENT_RED),
    ("Databases",     "H2, MySQL, PostgreSQL,\nMongoDB, Redis, Oracle, SQL Server", SPRING_GREEN),
    ("Frontend",      "AngularJS 1.2.16 (EOL)\nBootstrap 3.1.1", ACCENT_ORANGE),
    ("Deployment",    "Cloud Foundry\nsingle cf push", SPRING_GREEN),
    ("Observability", "Spring Actuator\n(all endpoints exposed)", ACCENT_ORANGE),
]

cols = [(0.4, 2.9), (4.55, 2.9), (9.0, 4.25)]
pos = [(0.4, 1.35), (4.55, 1.35), (9.0, 1.35)]

for idx, (label, value, clr) in enumerate(stack_items):
    col_idx = 0 if idx < 3 else (1 if idx < 6 else 2)
    row_idx = idx % 3
    x, w = [(0.4, 4.0), (4.55, 4.3), (9.0, 4.15)][col_idx]
    y = 1.35 + row_idx * 1.8

    card(sl, x, y, w, 1.6, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, x, y, 0.18, 1.6, fill=clr)
    add_text(sl, label.upper(), x + 0.28, y + 0.12, w - 0.35, 0.38,
             size=Pt(9), bold=True, color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(sl, value, x + 0.28, y + 0.48, w - 0.35, 0.95,
             size=Pt(13), bold=True, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Why Modernize
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Why Modernize?", "Pain points, risks, and business drivers")
bottom_bar(sl)

pain_points = [
    ("⚠  End-of-Life Stack",
     "Spring Boot 2.4 reached EOL Feb 2023. Java 8 support ends 2030 but misses 10+ years of language improvements. No security patches from the community."),
    ("⚠  Cannot Scale Independently",
     "A spike in album reads forces scaling the entire JAR — including the UI and metadata endpoints. No way to scale catalog independently."),
    ("⚠  Single Point of Failure",
     "A database connectivity issue, an OOM from /errors/fill-heap, or a misconfigured profile brings the entire application — including the UI — offline."),
    ("⚠  Deployment Coupling",
     "Any change — UI tweak, API fix, or new database support — requires rebuilding and redeploying the full monolith. All teams release together."),
    ("⚠  Incorrect Configuration",
     "application.yml specified 'ProgressDialect' for PostgreSQL (a Progress DB dialect). Hibernate dialects were also unnecessary in Hibernate 6."),
    ("⚠  No Error Handling",
     "Validation failures and unhandled exceptions returned raw Spring error pages. No structured error contract for consumers."),
]

for i, (title, body) in enumerate(pain_points):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.35 + row * 1.9
    card(sl, x, y, 6.2, 1.75, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, x, y, 6.2, 0.06, fill=ACCENT_RED)
    add_text(sl, title, x + 0.18, y + 0.15, 5.8, 0.45,
             size=Pt(12), bold=True, color=WHITE)
    add_text(sl, body, x + 0.18, y + 0.6, 5.8, 1.0,
             size=Pt(10.5), color=RGBColor(0xBB, 0xBB, 0xBB))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Target State: Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Target State", "Three independently deployable services")
bottom_bar(sl)

services = [
    ("album-catalog-service", "Port 8080",
     ["Album CRUD REST API", "Profile-based DB selection", "Seed data loading", "CORS + exception handling"],
     SPRING_GREEN, 0.4),
    ("app-info-service", "Port 8081",
     ["GET /appinfo — active profiles", "GET /service — CF bindings", "Reads VCAP_SERVICES via CfEnv", "Independent restart"],
     ACCENT_BLUE, 4.75),
    ("spring-music-ui", "Port 8082",
     ["Serves AngularJS SPA (static)", "Reverse-proxy → catalog & info", "No API logic, no DB dependency", "Browser talks to single host"],
     ACCENT_ORANGE, 9.1),
]

# Browser box at top centre
add_rect(sl, 5.4, 1.3, 2.5, 0.75, fill=ACCENT_BLUE)
add_text(sl, "Browser  (AngularJS SPA)", 5.45, 1.38, 2.4, 0.6,
         size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow down to UI
add_text(sl, "│  All requests\n▼", 6.5, 2.1, 0.8, 0.55,
         size=Pt(11), color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

for name, port, bullets, clr, x in services:
    card(sl, x, 2.7, 4.15, 4.45, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, x, 2.7, 4.15, 0.07, fill=clr)
    add_text(sl, name, x + 0.15, 2.82, 3.85, 0.5,
             size=Pt(13), bold=True, color=WHITE)
    add_text(sl, port, x + 0.15, 3.28, 3.85, 0.38,
             size=Pt(11), color=clr)
    add_bullet_box(sl, bullets, x + 0.15, 3.68, 3.8, 3.0,
                   size=Pt(11), color=RGBColor(0xCC, 0xCC, 0xCC))

# Proxy arrows
add_text(sl, "proxy /albums/**\n◄────────────►", 4.2, 4.4, 0.5, 0.7,
         size=Pt(8.5), color=RGBColor(0x88, 0x88, 0x88))
add_text(sl, "proxy /appinfo\n◄────────────►", 8.55, 4.4, 0.5, 0.7,
         size=Pt(8.5), color=RGBColor(0x88, 0x88, 0x88))

# DB note under catalog
add_rect(sl, 0.4, 7.1, 4.15, 0.05, fill=SPRING_GREEN)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Target State: Tech stack upgrade
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Target State", "Modernized technology stack")
bottom_bar(sl)

upgrades = [
    ("Java",          "8",                "17 LTS",            "Language features: records, sealed classes, text blocks"),
    ("Spring Boot",   "2.4.0  (EOL)",     "3.2.3",             "Jakarta EE 10, virtual threads, ProblemDetail, native image"),
    ("Gradle",        "6.7",              "8.7",               "Build cache, toolchains, modern plugins {} DSL"),
    ("JUnit",         "4",                "5 (Jupiter)",        "Parameterised tests, extensions, @Nested"),
    ("Namespace",     "javax.*",          "jakarta.*",          "Required by Spring Boot 3 / Jakarta EE 10"),
    ("MySQL Driver",  "mysql-connector-java", "mysql-connector-j", "Updated group ID + JDBC URL driver class"),
    ("Dialects",      "Explicit (& wrong)", "Auto-detected",   "Hibernate 6 detects dialect from JDBC metadata"),
    ("YAML profiles", "spring.profiles:", "spring.config.activate.on-profile:", "Mandatory syntax change in Spring Boot 2.4+"),
]

# Header row
hdr_y = 1.35
for x, w, lbl in [(0.4, 2.2, "Component"), (2.65, 2.5, "Before"), (5.2, 2.5, "After"), (7.75, 5.3, "Why it matters")]:
    add_rect(sl, x, hdr_y, w, 0.45, fill=SPRING_GREEN)
    add_text(sl, lbl, x + 0.1, hdr_y + 0.07, w - 0.15, 0.35,
             size=Pt(12), bold=True, color=WHITE)

for i, (comp, before, after, why) in enumerate(upgrades):
    row_y = 1.82 + i * 0.62
    bg = RGBColor(0x3A, 0x3A, 0x3A) if i % 2 == 0 else RGBColor(0x32, 0x32, 0x32)
    for x, w in [(0.4, 13.65)]:
        add_rect(sl, x, row_y, w, 0.6, fill=bg)
    add_text(sl, comp,   0.5,  row_y + 0.12, 2.1,  0.4, size=Pt(11), bold=True, color=WHITE)
    add_text(sl, before, 2.65, row_y + 0.12, 2.45, 0.4, size=Pt(11), color=ACCENT_RED)
    add_text(sl, after,  5.2,  row_y + 0.12, 2.45, 0.4, size=Pt(11), color=SPRING_GREEN)
    add_text(sl, why,    7.75, row_y + 0.12, 5.2,  0.4, size=Pt(10), color=RGBColor(0xBB, 0xBB, 0xBB))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Differences: Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Key Differences", "Architecture — As-Is vs To-Be")
bottom_bar(sl)

# Divider
add_rect(sl, 6.55, 1.35, 0.08, 5.8, fill=SPRING_GREEN)
add_text(sl, "AS-IS", 1.5, 1.35, 3.0, 0.5,
         size=Pt(18), bold=True, color=ACCENT_RED, align=PP_ALIGN.CENTER)
add_text(sl, "TO-BE", 8.5, 1.35, 3.0, 0.5,
         size=Pt(18), bold=True, color=SPRING_GREEN, align=PP_ALIGN.CENTER)

asis = [
    "Single deployable JAR (~1 GB allocation)",
    "All components share one JVM process",
    "One database binding per deployment",
    "UI + API + metadata in same package",
    "Full redeploy for any change",
    "Single failure point for all capabilities",
    "No CORS — browser & API on same origin only",
    "No structured exception handling",
]
tobe = [
    "Three independently deployable JARs",
    "Each service runs in its own JVM process",
    "Catalog DB fully isolated from other services",
    "UI / Catalog API / Metadata clearly separated",
    "Each service deploys independently",
    "Catalog failure does not affect UI or info",
    "CORS enabled — UI proxies to backend services",
    "GlobalExceptionHandler with RFC 7807 format",
]

for i, (a, b) in enumerate(zip(asis, tobe)):
    y = 1.92 + i * 0.65
    add_rect(sl, 0.3, y, 6.1, 0.57, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, 0.3, y, 0.15, 0.57, fill=ACCENT_RED)
    add_text(sl, a, 0.55, y + 0.1, 5.8, 0.38, size=Pt(11), color=WHITE)

    add_rect(sl, 6.75, y, 6.25, 0.57, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, 6.75, y, 0.15, 0.57, fill=SPRING_GREEN)
    add_text(sl, b, 7.0, y + 0.1, 5.9, 0.38, size=Pt(11), color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Key Differences: Code
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Key Differences", "Code & Configuration — As-Is vs To-Be")
bottom_bar(sl)

rows = [
    ("HTTP Status codes",    "200 for all operations",             "201 Created (add), 204 No Content (delete)"),
    ("Validation errors",    "Raw Spring whitelabel error page",   "RFC 7807 ProblemDetail with field details"),
    ("Logging",              "String concatenation in logger",     "SLF4J parameterised {} tokens"),
    ("ApplicationInfo",      "Mutable POJO, 28 lines",             "Java 17 record, 3 lines"),
    ("MySQL driver class",   "com.mysql.jdbc.Driver (deprecated)", "com.mysql.cj.jdbc.Driver"),
    ("PostgreSQL dialect",   "ProgressDialect (wrong DB!)",        "Removed — Hibernate 6 auto-detects"),
    ("YAML profile key",     "spring.profiles: mysql",             "spring.config.activate.on-profile: mysql"),
    ("Build repositories",   "jcenter() + mavenCentral()",        "mavenCentral() only (jcenter shutdown)"),
    ("Build script style",   "Legacy buildscript {} + apply",      "Modern plugins {} DSL"),
    ("IdentifierGenerator",  "SessionImplementor (Hibernate 5)",   "SharedSessionContractImplementor (Hibernate 6)"),
]

# Header
for x, w, lbl in [(0.35, 3.5, "Area"), (3.9, 4.55, "Before (As-Is)"), (8.5, 4.65, "After (To-Be)")]:
    add_rect(sl, x, 1.35, w, 0.42, fill=SPRING_GREEN)
    add_text(sl, lbl, x + 0.1, 1.38, w - 0.15, 0.35,
             size=Pt(12), bold=True, color=WHITE)

for i, (area, before, after) in enumerate(rows):
    y = 1.79 + i * 0.54
    bg = RGBColor(0x3A, 0x3A, 0x3A) if i % 2 == 0 else RGBColor(0x32, 0x32, 0x32)
    add_rect(sl, 0.35, y, 13.3, 0.52, fill=bg)
    add_text(sl, area,   0.45, y + 0.1, 3.35, 0.35, size=Pt(10.5), bold=True,  color=WHITE)
    add_text(sl, before, 3.9,  y + 0.1, 4.45, 0.35, size=Pt(10.5), color=ACCENT_RED)
    add_text(sl, after,  8.5,  y + 0.1, 4.55, 0.35, size=Pt(10.5), color=SPRING_GREEN)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Migration Summary: What changed / unchanged
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Migration Summary", "What changed vs what was preserved")
bottom_bar(sl)

changed = [
    "Spring Boot 2.4 → 3.2.3",
    "Java 8 → Java 17",
    "Gradle 6.7 → Gradle 8.7",
    "JUnit 4 → JUnit 5 (Jupiter)",
    "javax.* → jakarta.*",
    "mysql-connector-java → mysql-connector-j",
    "Removed explicit Hibernate dialects",
    "Fixed YAML profile syntax",
    "Added GlobalExceptionHandler",
    "Added CORS configuration (WebConfig)",
    "ApplicationInfo: POJO → Java record",
    "AlbumController: @RequestMapping → @GetMapping etc.",
    "HTTP 200 → 201 Created / 204 No Content",
    "Hibernate 6 IdentifierGenerator API",
    "Removed jcenter(), eclipse-wtp plugin",
    "Decomposed into 3 services",
    "Added ProxyController in spring-music-ui",
]

unchanged = [
    "All AngularJS JS files (albums.js, app.js …)",
    "All HTML templates",
    "All CSS stylesheets",
    "albums.json seed data",
    "Profile-based DB selection logic",
    "AlbumRepositoryPopulator behaviour",
    "JpaAlbumRepository (interface only)",
    "MongoAlbumRepository (interface only)",
    "RedisAlbumRepository core logic",
    "SpringApplicationContextInitializer algorithm",
    "CF deployment model (cf push per service)",
    "Actuator endpoints exposure",
]

# Changed column
add_rect(sl, 0.35, 1.35, 6.1, 0.45, fill=ACCENT_RED)
add_text(sl, "CHANGED", 0.45, 1.38, 5.9, 0.38,
         size=Pt(13), bold=True, color=WHITE)
add_bullet_box(sl, changed,
               0.4, 1.85, 5.9, 5.3,
               size=Pt(10.5), color=RGBColor(0xCC, 0xCC, 0xCC))

# Unchanged column
add_rect(sl, 6.85, 1.35, 6.1, 0.45, fill=SPRING_GREEN)
add_text(sl, "UNCHANGED", 6.95, 1.38, 5.9, 0.38,
         size=Pt(13), bold=True, color=WHITE)
add_bullet_box(sl, unchanged,
               6.9, 1.85, 5.9, 5.3,
               size=Pt(10.5), color=RGBColor(0xCC, 0xCC, 0xCC))

# Divider
add_rect(sl, 6.6, 1.35, 0.08, 5.8, fill=RGBColor(0x55, 0x55, 0x55))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Next Steps
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
green_bar(sl, "Next Steps & Known Gaps")
bottom_bar(sl)

next_steps = [
    ("🔐  Authentication & Authorization",
     "Spring Security with OAuth2 resource server.\nCurrently absent — all endpoints publicly accessible.",
     ACCENT_RED),
    ("📡  Distributed Tracing",
     "Add Micrometer Tracing. Propagate correlation ID\nthrough ProxyController for end-to-end visibility.",
     ACCENT_ORANGE),
    ("🖥  Frontend Modernization",
     "AngularJS 1.2.16 reached EOL in 2021.\nRecommend migrating to Angular 17, React, or Vue.",
     ACCENT_ORANGE),
    ("🔎  Service Discovery",
     "Backend URLs in spring-music-ui are static config.\nAdd Spring Cloud Discovery for production CF routing.",
     ACCENT_BLUE),
    ("⚙  @UuidGenerator Migration",
     "@GenericGenerator deprecated in Hibernate 6.\nReplace with @UuidGenerator once dash-free format is unneeded.",
     ACCENT_BLUE),
    ("📦  Container / Kubernetes",
     "Services are now independently containerisable.\nNext logical step: Dockerfile + Kubernetes Helm charts.",
     SPRING_GREEN),
]

for i, (title, body, clr) in enumerate(next_steps):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.35 + row * 1.95
    card(sl, x, y, 6.2, 1.8, fill=RGBColor(0x3A, 0x3A, 0x3A))
    add_rect(sl, x, y, 0.18, 1.8, fill=clr)
    add_text(sl, title, x + 0.28, y + 0.15, 5.8, 0.48,
             size=Pt(12), bold=True, color=WHITE)
    add_text(sl, body,  x + 0.28, y + 0.65, 5.8, 1.0,
             size=Pt(10.5), color=RGBColor(0xBB, 0xBB, 0xBB))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Thank You / Questions
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, 0, 0, 0.35, 7.5, fill=SPRING_GREEN)
add_rect(sl, 0.35, 3.3, 12.98, 0.06, fill=SPRING_GREEN)

add_text(sl, "Thank You", 0.7, 1.2, 11.5, 1.2,
         size=Pt(52), bold=True, color=WHITE)
add_text(sl, "Questions?", 0.7, 2.45, 11.5, 0.75,
         size=Pt(30), color=SPRING_GREEN)

summary_lines = [
    "✓  Spring Boot 2.4  →  3.2.3   |   Java 8  →  17   |   Gradle 6.7  →  8.7",
    "✓  Monolith decomposed into 3 independent services",
    "✓  javax.*  →  jakarta.*   |   Deprecated APIs replaced   |   RFC 7807 error handling added",
    "✓  AngularJS SPA unchanged — ProxyController preserves all relative URLs",
]
add_bullet_box(sl, summary_lines, 0.7, 3.6, 12.0, 3.2,
               size=Pt(13), color=RGBColor(0xCC, 0xCC, 0xCC), bullet="")

add_text(sl, "Source: spring-music/  →  spring-music-new/", 0.7, 6.8, 9.0, 0.4,
         size=Pt(11), color=RGBColor(0x77, 0x77, 0x77))


# ─── Save ────────────────────────────────────────────────────────────────────
output_path = "spring_music_modernization.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
